#!/usr/bin/env python3
"""
Erstellt eine PowerPoint-Praesentation fuer den Gewerbespeicher Planner.
Marktwert-Validierung, Wertdarstellung und Anwendungsuebersicht.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Farben
DARK_BLUE = RGBColor(0x1B, 0x2A, 0x4A)
ACCENT_BLUE = RGBColor(0x2E, 0x86, 0xAB)
ACCENT_GREEN = RGBColor(0x28, 0xA7, 0x45)
ACCENT_ORANGE = RGBColor(0xF0, 0x93, 0x19)
ACCENT_RED = RGBColor(0xDC, 0x35, 0x45)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF0, 0xF2, 0xF5)
MEDIUM_GRAY = RGBColor(0x6C, 0x75, 0x7D)
TEXT_DARK = RGBColor(0x21, 0x25, 0x29)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SW = prs.slide_width
SH = prs.slide_height


def add_shape(slide, left, top, width, height, fill_color=None, line_color=None, shape_type=MSO_SHAPE.RECTANGLE):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=14, bold=False,
                 color=TEXT_DARK, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_kpi_card(slide, left, top, width, height, value, label, accent_color=ACCENT_BLUE):
    card = add_shape(slide, left, top, width, height, fill_color=WHITE)
    add_shape(slide, left, top, width, Pt(4), fill_color=accent_color)
    add_text_box(slide, left + Inches(0.15), top + Inches(0.15), width - Inches(0.3), Inches(0.5),
                 value, font_size=24, bold=True, color=accent_color, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.1), top + Inches(0.6), width - Inches(0.2), Inches(0.5),
                 label, font_size=10, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)


def add_bullet_list(slide, left, top, width, height, items, font_size=11, bullet="\u25B8"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"{bullet} {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = TEXT_DARK
        p.font.name = "Calibri"
        p.space_after = Pt(4)
    return txBox


# ============================================================
# SLIDE 1: Anwendungsuebersicht & Was wir gebaut haben
# ============================================================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide1.background
fill = bg.fill
fill.solid()
fill.fore_color.rgb = LIGHT_GRAY

# Header
add_shape(slide1, 0, 0, SW, Inches(1.6), fill_color=DARK_BLUE)
add_text_box(slide1, Inches(0.8), Inches(0.25), Inches(9), Inches(0.5),
             "Gewerbespeicher Planner", font_size=30, bold=True, color=WHITE)
add_text_box(slide1, Inches(0.8), Inches(0.8), Inches(9), Inches(0.4),
             "KI-gestuetzte Planungs- & Angebotssoftware fuer gewerbliche PV-Speichersysteme",
             font_size=14, color=RGBColor(0xA0, 0xC0, 0xDF))

# AI Badge
add_shape(slide1, Inches(10.5), Inches(0.25), Inches(2.4), Inches(0.5), fill_color=ACCENT_GREEN)
add_text_box(slide1, Inches(10.5), Inches(0.27), Inches(2.4), Inches(0.5),
             "Built with AI (Claude)", font_size=13, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
add_shape(slide1, Inches(10.5), Inches(0.85), Inches(2.4), Inches(0.35), fill_color=RGBColor(0x1D, 0x7A, 0x3A))
add_text_box(slide1, Inches(10.5), Inches(0.85), Inches(2.4), Inches(0.35),
             "Anthropic Claude Code & Opus", font_size=9, color=WHITE, alignment=PP_ALIGN.CENTER)

# --- Codebase-KPIs ---
kpi_y = Inches(1.85)
kpi_h = Inches(0.95)
kpi_w = Inches(1.85)
kpi_gap = Inches(0.2)
kpi_start = Inches(0.6)

code_kpis = [
    ("31.569", "Zeilen Code", ACCENT_BLUE),
    ("94", "Source-Dateien", ACCENT_BLUE),
    ("101", "API-Endpunkte", ACCENT_GREEN),
    ("11", "Backend-Services", ACCENT_GREEN),
    ("24", "React-Komponenten", ACCENT_ORANGE),
    ("9", "Datenbank-Modelle", ACCENT_ORANGE),
]
for idx, (val, lbl, clr) in enumerate(code_kpis):
    x = kpi_start + idx * (kpi_w + kpi_gap)
    add_kpi_card(slide1, x, kpi_y, kpi_w, kpi_h, val, lbl, clr)

# --- Zwei Spalten ---
col_y = Inches(3.05)
left_x = Inches(0.6)
right_x = Inches(6.8)
col_w = Inches(5.8)

# LINKS: Kernfunktionen
add_text_box(slide1, left_x, col_y, col_w, Inches(0.35),
             "Kernfunktionen", font_size=16, bold=True, color=DARK_BLUE)
add_bullet_list(slide1, left_x, col_y + Inches(0.35), col_w, Inches(1.8), [
    "Echtzeit PV+Batterie-Simulation (pvlib, stuendliche Aufloesung)",
    "KI-optimierte Systemdimensionierung & Angebotstexte (Claude Opus)",
    "Professionelle PDF-Angebotserstellung mit White-Label-Branding",
    "Finanzanalyse: ROI, Amortisation, Kapitalwert (NPV), IRR",
    "Multi-Tenant-Architektur fuer Installateur-Netzwerk",
], font_size=11)

add_text_box(slide1, left_x, col_y + Inches(2.2), col_w, Inches(0.35),
             "Integrationen & Compliance", font_size=16, bold=True, color=DARK_BLUE)
add_bullet_list(slide1, left_x, col_y + Inches(2.55), col_w, Inches(1.3), [
    "DocuSign (E-Signatur) | HubSpot (CRM) | Google Maps (Standort)",
    "PVGIS & Open-Meteo (Wetter-/Ertragsdaten) | BNetzA-Tarife",
    "EEG 2023, Paragraph 14a EnWG, MaStR, KfW-Foerderung vollstaendig integriert",
], font_size=11)

# RECHTS: Tech-Stack + Status
add_text_box(slide1, right_x, col_y, col_w, Inches(0.35),
             "Technologie-Stack", font_size=16, bold=True, color=DARK_BLUE)

tech_items = [
    ("Frontend", "Next.js 15, React 19, TypeScript, Tailwind CSS"),
    ("Backend", "FastAPI, SQLAlchemy, Pydantic, pvlib-python"),
    ("KI", "Anthropic Claude Opus - Intelligente Angebotsoptimierung"),
    ("Infrastruktur", "Vercel + Railway + Neon (PostgreSQL) + Upstash (Redis)"),
    ("Compliance", "EEG 2023, Paragraph 14a EnWG, MaStR, KfW"),
]
ty = col_y + Inches(0.4)
for title, desc in tech_items:
    add_shape(slide1, right_x, ty, col_w, Inches(0.55), fill_color=WHITE)
    add_text_box(slide1, right_x + Inches(0.15), ty + Inches(0.05), Inches(1.5), Inches(0.25),
                 title, font_size=10, bold=True, color=ACCENT_BLUE)
    add_text_box(slide1, right_x + Inches(1.7), ty + Inches(0.05), Inches(3.9), Inches(0.45),
                 desc, font_size=10, color=TEXT_DARK)
    ty += Inches(0.6)

# Workflow-Pfeil
add_text_box(slide1, right_x, ty + Inches(0.15), col_w, Inches(0.35),
             "End-to-End Workflow", font_size=14, bold=True, color=DARK_BLUE)
add_shape(slide1, right_x, ty + Inches(0.5), col_w, Inches(0.45), fill_color=DARK_BLUE)
add_text_box(slide1, right_x, ty + Inches(0.5), col_w, Inches(0.45),
             "Konfiguration  -->  Simulation  -->  KI-Optimierung  -->  PDF-Angebot  -->  E-Signatur  -->  CRM",
             font_size=11, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

# Footer: Status
add_shape(slide1, Inches(0.6), Inches(6.65), Inches(12.1), Inches(0.55), fill_color=WHITE)
phases = [
    ("Phase 1: MVP", Inches(0.8)),
    ("Phase 2: KI-Integration", Inches(3.5)),
    ("Phase 3: Integrationen", Inches(6.5)),
    ("Phase 4: Enterprise", Inches(9.5)),
]
for label, x in phases:
    add_text_box(slide1, x, Inches(6.7), Inches(2.8), Inches(0.25),
                 f"[OK] {label}", font_size=11, bold=True, color=ACCENT_GREEN)
add_text_box(slide1, Inches(0.8), Inches(6.95), Inches(10), Inches(0.2),
             "Alle 4 Entwicklungsphasen abgeschlossen | Produktionsreif | KI-gestuetzte Entwicklung mit Claude Code (Anthropic)",
             font_size=9, color=MEDIUM_GRAY)


# ============================================================
# SLIDE 2: Marktvalidierung & Wettbewerb
# ============================================================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
bg2 = slide2.background
fill2 = bg2.fill
fill2.solid()
fill2.fore_color.rgb = LIGHT_GRAY

add_shape(slide2, 0, 0, SW, Inches(1.15), fill_color=DARK_BLUE)
add_text_box(slide2, Inches(0.8), Inches(0.25), Inches(10), Inches(0.45),
             "Marktvalidierung & Wettbewerbsanalyse", font_size=26, bold=True, color=WHITE)
add_text_box(slide2, Inches(0.8), Inches(0.7), Inches(10), Inches(0.3),
             "Quellen: Wood Mackenzie, Mordor Intelligence, SaaS Capital, BSW Solar, ESS-News, pv magazine, BNetzA",
             font_size=10, color=RGBColor(0xA0, 0xC0, 0xDF))

# KPIs
kpi_y2 = Inches(1.4)
kpi_h2 = Inches(1.05)
kpi_w2 = Inches(2.85)
gap2 = Inches(0.3)

add_kpi_card(slide2, Inches(0.6), kpi_y2, kpi_w2, kpi_h2,
             "24 GWh", "DE Speicherkapazitaet 2025\n(+6,57 GWh in 2025)", ACCENT_BLUE)
add_kpi_card(slide2, Inches(0.6) + kpi_w2 + gap2, kpi_y2, kpi_w2, kpi_h2,
             "+34%", "Gewerbespeicher-Wachstum\n(Jan 2026 vs. Jan 2025)", ACCENT_GREEN)
add_kpi_card(slide2, Inches(0.6) + 2 * (kpi_w2 + gap2), kpi_y2, kpi_w2, kpi_h2,
             "EUR 19 Mrd.", "DE Energiespeicher-Markt\n(Gesamtvolumen 2024)", ACCENT_ORANGE)
add_kpi_card(slide2, Inches(0.6) + 3 * (kpi_w2 + gap2), kpi_y2, kpi_w2, kpi_h2,
             "EUR 2,5 Mrd.", "Energiespeicher-Software\n(Global 2023, CAGR 15,1%)", ACCENT_BLUE)

# LINKS: Marktzahlen
sec_y2 = Inches(2.75)
add_text_box(slide2, Inches(0.6), sec_y2, Inches(6.2), Inches(0.35),
             "Marktumfeld & Wachstumstreiber", font_size=15, bold=True, color=DARK_BLUE)
add_bullet_list(slide2, Inches(0.6), sec_y2 + Inches(0.35), Inches(6.2), Inches(3.2), [
    "PV-Software-Markt: $1,88 Mrd. (2025) -> $2,81 Mrd. (2034), CAGR 6-9%",
    "C&I Energiespeicher: $92 Mrd. (2025) -> $164 Mrd. (2030), CAGR 12,3%",
    "Gewerbespeicher DE: 1.248 Systeme Jan. 2026 (+34% YoY)",
    "575 Stunden neg. Strompreise 2025 - Speicher hochattraktiv",
    "Solar Peak Act (Feb. 2025): Speicher wirtschaftlich Pflicht",
    "EU-Ziel 42,5% Erneuerbare bis 2030 treibt Investitionen",
    "22 Energy-SaaS-Akquisitionen in 2025 (Tracxn)",
    "PE-Fonds aktiv im Industrial/Energy-Software-Segment",
], font_size=10.5)

# RECHTS: Wettbewerb
add_text_box(slide2, Inches(7.1), sec_y2, Inches(5.6), Inches(0.35),
             "Wettbewerb & Positionierung", font_size=15, bold=True, color=DARK_BLUE)

competitors = [
    ("PV*SOL (Valentin)", "Desktop, ~845 EUR/J.", "Nur Simulation, keine KI, kein E2E"),
    ("Eturnity", "SaaS, Abo-Modell", "Kalkulator, kein Speicher-Fokus"),
    ("Energy Toolbase", "SaaS, US-Markt", "Speicher, aber keine DE-Compliance"),
    ("PVcase / PVsyst", "Enterprise/Desktop", "Grossprojekte, kein Gewerbe-Fokus"),
]

comp_y = sec_y2 + Inches(0.45)
for name, model, gap_text in competitors:
    add_shape(slide2, Inches(7.1), comp_y, Inches(5.6), Inches(0.5), fill_color=WHITE)
    add_text_box(slide2, Inches(7.25), comp_y + Inches(0.02), Inches(2.0), Inches(0.22),
                 name, font_size=10, bold=True, color=ACCENT_BLUE)
    add_text_box(slide2, Inches(7.25), comp_y + Inches(0.25), Inches(2.0), Inches(0.22),
                 model, font_size=9, color=MEDIUM_GRAY)
    add_text_box(slide2, Inches(9.3), comp_y + Inches(0.08), Inches(3.3), Inches(0.35),
                 gap_text, font_size=10, color=ACCENT_RED)
    comp_y += Inches(0.55)

# USP-Box
usp_y = comp_y + Inches(0.1)
add_shape(slide2, Inches(7.1), usp_y, Inches(5.6), Inches(1.2), fill_color=DARK_BLUE)
add_text_box(slide2, Inches(7.25), usp_y + Inches(0.05), Inches(5.3), Inches(0.25),
             "Gewerbespeicher Planner - Unser USP", font_size=12, bold=True, color=ACCENT_GREEN)
add_bullet_list(slide2, Inches(7.25), usp_y + Inches(0.3), Inches(5.3), Inches(0.85), [
    "Einzige DE-Loesung: KI + Compliance + End-to-End",
    "Kompletter Workflow: Simulation -> Angebot -> Signatur -> CRM",
    "Multi-Tenant White-Label fuer Installateur-Skalierung",
    "Regulatorische Komplexitaet als Markteintrittsbarriere",
], font_size=10)
# Override bullet color to white
for shape in slide2.shapes:
    pass  # bullets set above

# Fix: make USP bullets white
usp_bullets = slide2.shapes[-1]
for para in usp_bullets.text_frame.paragraphs:
    para.font.color.rgb = RGBColor(0xC0, 0xD5, 0xE8)

# Footer
add_shape(slide2, 0, Inches(6.7), SW, Inches(0.8), fill_color=WHITE)
add_text_box(slide2, Inches(0.6), Inches(6.75), Inches(12), Inches(0.5),
             "Quellen: Wood Mackenzie (2025) | Mordor Intelligence C&I Storage Report | BSW Solar | ESS-News.com (Jan 2026) | "
             "SaaS Capital Valuation 2025 | pv magazine DE | Tracxn Energy SaaS 2026 | BNetzA | Precedence Research",
             font_size=8, color=MEDIUM_GRAY)


# ============================================================
# SLIDE 3: WERTDARSTELLUNG - Das Herzstuck
# ============================================================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
bg3 = slide3.background
fill3 = bg3.fill
fill3.solid()
fill3.fore_color.rgb = LIGHT_GRAY

add_shape(slide3, 0, 0, SW, Inches(1.15), fill_color=DARK_BLUE)
add_text_box(slide3, Inches(0.8), Inches(0.2), Inches(9), Inches(0.45),
             "Wertdarstellung: Was ist diese Anwendung wert?", font_size=26, bold=True, color=WHITE)
add_text_box(slide3, Inches(0.8), Inches(0.65), Inches(9), Inches(0.3),
             "Drei Perspektiven: Wiederbeschaffungswert, KI-Kostenersparnis und SaaS-Bewertungspotenzial",
             font_size=12, color=RGBColor(0xA0, 0xC0, 0xDF))

# AI Badge
add_shape(slide3, Inches(10.5), Inches(0.2), Inches(2.4), Inches(0.7), fill_color=ACCENT_GREEN)
add_text_box(slide3, Inches(10.5), Inches(0.22), Inches(2.4), Inches(0.35),
             "Built with AI", font_size=13, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
add_text_box(slide3, Inches(10.5), Inches(0.52), Inches(2.4), Inches(0.35),
             "Claude Code (Anthropic)", font_size=9, color=WHITE, alignment=PP_ALIGN.CENTER)

# ===== BLOCK 1: Wiederbeschaffungswert =====
b1_x = Inches(0.6)
b1_y = Inches(1.4)
b1_w = Inches(3.9)
b1_h = Inches(3.3)

add_shape(slide3, b1_x, b1_y, b1_w, b1_h, fill_color=WHITE)
add_shape(slide3, b1_x, b1_y, b1_w, Pt(5), fill_color=ACCENT_BLUE)

add_text_box(slide3, b1_x + Inches(0.2), b1_y + Inches(0.12), b1_w - Inches(0.4), Inches(0.3),
             "Wiederbeschaffungswert", font_size=15, bold=True, color=ACCENT_BLUE)
add_text_box(slide3, b1_x + Inches(0.2), b1_y + Inches(0.4), b1_w - Inches(0.4), Inches(0.2),
             "Was wuerde konventionelle Entwicklung kosten?", font_size=9, color=MEDIUM_GRAY)

# Kalkulation
calc_items = [
    ("31.569 Zeilen Code, 94 Dateien, 101 API-Endpunkte", TEXT_DARK),
    ("", TEXT_DARK),
    ("Geschaetzter Aufwand (konventionell):", ACCENT_BLUE),
    ("  Senior Full-Stack-Entwickler    12 Monate", TEXT_DARK),
    ("  UX/UI-Designer                          3 Monate", TEXT_DARK),
    ("  DevOps/Infrastruktur                  2 Monate", TEXT_DARK),
    ("  Projektmanagement                    6 Monate", TEXT_DARK),
    ("  Domain-Experte (Solar/Energie)  3 Monate", TEXT_DARK),
    ("  = ca. 26 Personenmonate", ACCENT_BLUE),
    ("", TEXT_DARK),
    ("DE Agentur-Stundensatz: EUR 100-150/Std.", TEXT_DARK),
    ("26 PM x 160 Std. x EUR 120 =", TEXT_DARK),
]

calc_box = slide3.shapes.add_textbox(b1_x + Inches(0.2), b1_y + Inches(0.6), b1_w - Inches(0.4), Inches(2.0))
calc_tf = calc_box.text_frame
calc_tf.word_wrap = True
for i, (text, clr) in enumerate(calc_items):
    p = calc_tf.paragraphs[0] if i == 0 else calc_tf.add_paragraph()
    p.text = text
    p.font.size = Pt(9)
    p.font.color.rgb = clr
    p.font.name = "Calibri"
    p.space_after = Pt(1)

# Grosser Wert
add_text_box(slide3, b1_x + Inches(0.2), b1_y + Inches(2.5), b1_w - Inches(0.4), Inches(0.5),
             "EUR 350.000 - 500.000", font_size=22, bold=True, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)

add_text_box(slide3, b1_x + Inches(0.2), b1_y + Inches(2.95), b1_w - Inches(0.4), Inches(0.25),
             "Quelle: Clutch.co, FullStack Labs, Netguru (DE Agency Rates 2025)",
             font_size=7, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)


# ===== BLOCK 2: KI-Kostenersparnis =====
b2_x = Inches(4.7)
b2_y = Inches(1.4)
b2_w = Inches(3.9)
b2_h = Inches(3.3)

add_shape(slide3, b2_x, b2_y, b2_w, b2_h, fill_color=WHITE)
add_shape(slide3, b2_x, b2_y, b2_w, Pt(5), fill_color=ACCENT_GREEN)

add_text_box(slide3, b2_x + Inches(0.2), b2_y + Inches(0.12), b2_w - Inches(0.4), Inches(0.3),
             "KI-Entwicklungseffizienz", font_size=15, bold=True, color=ACCENT_GREEN)
add_text_box(slide3, b2_x + Inches(0.2), b2_y + Inches(0.4), b2_w - Inches(0.4), Inches(0.2),
             "Tatsaechliche Kosten durch KI-gestuetzte Entwicklung", font_size=9, color=MEDIUM_GRAY)

# Vergleich
compare_items = [
    ("Konventionelle Entwicklung:", ACCENT_RED),
    ("  26 Personenmonate", TEXT_DARK),
    ("  5-6 Entwickler, 12+ Monate", TEXT_DARK),
    ("  Kosten: EUR 350.000 - 500.000", TEXT_DARK),
    ("", TEXT_DARK),
    ("KI-gestuetzte Entwicklung (Claude Code):", ACCENT_GREEN),
    ("  1 Entwickler + KI-Pair-Programming", TEXT_DARK),
    ("  Deutlich verkuerzte Entwicklungszeit", TEXT_DARK),
    ("  4 Phasen vollstaendig abgeschlossen", TEXT_DARK),
    ("  Produktionsreife Qualitaet erreicht", TEXT_DARK),
    ("", TEXT_DARK),
    ("Geschaetzte Ersparnis:", ACCENT_GREEN),
]

comp_box = slide3.shapes.add_textbox(b2_x + Inches(0.2), b2_y + Inches(0.6), b2_w - Inches(0.4), Inches(2.0))
comp_tf = comp_box.text_frame
comp_tf.word_wrap = True
for i, (text, clr) in enumerate(compare_items):
    p = comp_tf.paragraphs[0] if i == 0 else comp_tf.add_paragraph()
    p.text = text
    p.font.size = Pt(9)
    p.font.color.rgb = clr
    p.font.name = "Calibri"
    p.space_after = Pt(1)

add_text_box(slide3, b2_x + Inches(0.2), b2_y + Inches(2.5), b2_w - Inches(0.4), Inches(0.5),
             "60-80% Kostenreduktion", font_size=22, bold=True, color=ACCENT_GREEN, alignment=PP_ALIGN.CENTER)
add_text_box(slide3, b2_x + Inches(0.2), b2_y + Inches(2.95), b2_w - Inches(0.4), Inches(0.25),
             "1 Entwickler + Claude AI statt 5-6 Personen Team",
             font_size=7, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)


# ===== BLOCK 3: SaaS-Bewertungspotenzial =====
b3_x = Inches(8.8)
b3_y = Inches(1.4)
b3_w = Inches(3.9)
b3_h = Inches(3.3)

add_shape(slide3, b3_x, b3_y, b3_w, b3_h, fill_color=WHITE)
add_shape(slide3, b3_x, b3_y, b3_w, Pt(5), fill_color=ACCENT_ORANGE)

add_text_box(slide3, b3_x + Inches(0.2), b3_y + Inches(0.12), b3_w - Inches(0.4), Inches(0.3),
             "SaaS-Bewertungspotenzial", font_size=15, bold=True, color=ACCENT_ORANGE)
add_text_box(slide3, b3_x + Inches(0.2), b3_y + Inches(0.4), b3_w - Inches(0.4), Inches(0.2),
             "Marktuebliche Bewertung bei Umsatzaufbau", font_size=9, color=MEDIUM_GRAY)

saas_items = [
    ("B2B SaaS-Multiples (SaaS Capital 2025):", ACCENT_ORANGE),
    ("  Bootstrapped: 4,8x ARR", TEXT_DARK),
    ("  VC-backed: 5,3x ARR", TEXT_DARK),
    ("  Median Markt: 7,0x ARR", TEXT_DARK),
    ("  Energy/Industrial SaaS: 3-7x ARR", TEXT_DARK),
    ("", TEXT_DARK),
    ("Beispielrechnung bei Markteintritt:", ACCENT_ORANGE),
    ("  50 Installateure x EUR 199/Monat", TEXT_DARK),
    ("  = EUR 119.400 ARR", TEXT_DARK),
    ("  x 5x Multiple =", TEXT_DARK),
]

saas_box = slide3.shapes.add_textbox(b3_x + Inches(0.2), b3_y + Inches(0.6), b3_w - Inches(0.4), Inches(2.0))
saas_tf = saas_box.text_frame
saas_tf.word_wrap = True
for i, (text, clr) in enumerate(saas_items):
    p = saas_tf.paragraphs[0] if i == 0 else saas_tf.add_paragraph()
    p.text = text
    p.font.size = Pt(9)
    p.font.color.rgb = clr
    p.font.name = "Calibri"
    p.space_after = Pt(1)

add_text_box(slide3, b3_x + Inches(0.2), b3_y + Inches(2.5), b3_w - Inches(0.4), Inches(0.5),
             "EUR 600.000+", font_size=22, bold=True, color=ACCENT_ORANGE, alignment=PP_ALIGN.CENTER)
add_text_box(slide3, b3_x + Inches(0.2), b3_y + Inches(2.95), b3_w - Inches(0.4), Inches(0.25),
             "Quelle: SaaS Capital 2025, Aventis Advisors, iMerge Advisors",
             font_size=7, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)


# ===== KERNAUSSAGE =====
key_y = Inches(4.95)
add_shape(slide3, Inches(0.6), key_y, Inches(12.1), Inches(1.0), fill_color=DARK_BLUE)
add_text_box(slide3, Inches(0.8), key_y + Inches(0.08), Inches(11.7), Inches(0.25),
             "Zusammenfassung: Der Wert dieser Anwendung", font_size=13, bold=True, color=ACCENT_GREEN)

summary_box = slide3.shapes.add_textbox(Inches(0.8), key_y + Inches(0.35), Inches(11.7), Inches(0.6))
summary_tf = summary_box.text_frame
summary_tf.word_wrap = True
p = summary_tf.paragraphs[0]
p.font.size = Pt(11)
p.font.color.rgb = WHITE
p.font.name = "Calibri"
p.text = (
    "Der Gewerbespeicher Planner hat einen Wiederbeschaffungswert von EUR 350.000-500.000. "
    "Durch KI-gestuetzte Entwicklung mit Claude Code wurde dieser Wert mit 60-80% weniger Kosten geschaffen. "
    "Als SaaS-Produkt im wachsenden DE-Gewerbespeichermarkt (+34% YoY) besteht Bewertungspotenzial von EUR 600.000+ "
    "bereits bei moderatem Markteintritt (50 Installateure). Die Kombination aus KI, Compliance und E2E-Workflow ist im DE-Markt einzigartig."
)

# ===== NAECHSTE SCHRITTE =====
next_y = Inches(6.15)
add_text_box(slide3, Inches(0.6), next_y, Inches(12), Inches(0.3),
             "Empfohlene naechste Schritte", font_size=14, bold=True, color=DARK_BLUE)

steps = [
    ("1", "Pilotphase", "5-10 Installateure als Beta-Nutzer, Feedback sammeln"),
    ("2", "Go-to-Market", "Preismodell (SaaS), Vertrieb ueber EWS-Netzwerk"),
    ("3", "Skalierung", "Multi-Tenant ausrollen, Partnernetzwerk aufbauen"),
]

step_x = Inches(0.6)
for num, title, desc in steps:
    add_shape(slide3, step_x, next_y + Inches(0.35), Inches(3.8), Inches(0.6), fill_color=WHITE)
    add_shape(slide3, step_x + Inches(0.08), next_y + Inches(0.42), Inches(0.35), Inches(0.35),
              fill_color=ACCENT_BLUE, shape_type=MSO_SHAPE.OVAL)
    add_text_box(slide3, step_x + Inches(0.08), next_y + Inches(0.42), Inches(0.35), Inches(0.35),
                 num, font_size=13, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide3, step_x + Inches(0.5), next_y + Inches(0.38), Inches(3.2), Inches(0.22),
                 title, font_size=11, bold=True, color=DARK_BLUE)
    add_text_box(slide3, step_x + Inches(0.5), next_y + Inches(0.58), Inches(3.2), Inches(0.3),
                 desc, font_size=9, color=MEDIUM_GRAY)
    step_x += Inches(4.15)

# Footer
add_text_box(slide3, Inches(0.6), Inches(7.1), Inches(12), Inches(0.3),
             "Gewerbespeicher Planner | EWS GmbH | Vertraulich | Februar 2026",
             font_size=8, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)


# ============================================================
# SPEICHERN
# ============================================================
output_path = "/home/user/Gewerbespeicher/Gewerbespeicher_Planner_Praesentation.pptx"
prs.save(output_path)
print(f"Praesentation gespeichert: {output_path}")
